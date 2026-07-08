"""
Genera la presentación PPTX del proyecto.
Salida: docs/Propuesta_Deteccion_Armas.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

OUT = Path(__file__).parent / "Propuesta_Deteccion_Armas.pptx"

# Paleta — Midnight Executive
NAVY      = RGBColor(0x0D, 0x2C, 0x54)   # primario
ICE_BLUE  = RGBColor(0xCA, 0xDC, 0xFC)   # secundario
ACCENT    = RGBColor(0x1F, 0x6F, 0xEB)   # accent
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF4, 0xF6, 0xF8)
MUTED     = RGBColor(0x5A, 0x6C, 0x82)
DARK_TXT  = RGBColor(0x1A, 0x1F, 0x29)

# Fuentes
F_TITLE = "Calibri"
F_BODY  = "Calibri"


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, *,
                 size=14, bold=False, color=DARK_TXT, align=PP_ALIGN.LEFT,
                 font=F_BODY, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_filled_rect(slide, left, top, width, height, color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    shape.shadow.inherit = False
    return shape


def add_accent_bar(slide, top=Inches(0.55), height=Inches(0.06), color=ACCENT):
    """Barra delgada de acento (NO debajo del título — la usamos como detalle lateral)."""
    return add_filled_rect(slide, Inches(0.4), top, Inches(0.4), height, color)


def add_footer(slide, idx, total):
    add_text_box(
        slide, Inches(0.4), Inches(7.05), Inches(8), Inches(0.3),
        "Propuesta · Detección de Armas con IA",
        size=9, color=MUTED, align=PP_ALIGN.LEFT,
    )
    add_text_box(
        slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
        f"{idx} / {total}",
        size=9, color=MUTED, align=PP_ALIGN.RIGHT,
    )


def add_page_title(slide, title, subtitle=None):
    # Pequeña etiqueta lateral (motivo visual repetido)
    add_filled_rect(slide, Inches(0.4), Inches(0.6), Inches(0.08), Inches(0.55), ACCENT)
    add_text_box(
        slide, Inches(0.65), Inches(0.55), Inches(11), Inches(0.65),
        title, size=30, bold=True, color=NAVY, font=F_TITLE,
    )
    if subtitle:
        add_text_box(
            slide, Inches(0.65), Inches(1.15), Inches(11), Inches(0.35),
            subtitle, size=14, color=MUTED, font=F_BODY,
        )


# ───────── Crear presentación ─────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

TOTAL_SLIDES = 16

# ───────── SLIDE 1: Portada ─────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, NAVY)

# Acento decorativo lateral
add_filled_rect(s, Inches(0), Inches(0), Inches(0.3), Inches(7.5), ACCENT)
# Marca esquina inferior derecha
add_filled_rect(s, Inches(10.5), Inches(6.8), Inches(2.5), Inches(0.08), ACCENT)

# Etiqueta superior
add_text_box(
    s, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
    "PROPUESTA · PROYECTO ACADÉMICO 2025–2026",
    size=14, bold=True, color=ICE_BLUE, font=F_TITLE,
)

# Título grande
add_text_box(
    s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.6),
    "Sistema de Detección\nde Armas con IA",
    size=54, bold=True, color=WHITE, font=F_TITLE, line_spacing=1.05,
)

# Subtítulo
add_text_box(
    s, Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.7),
    "Visión por computadora en tiempo real sobre Raspberry Pi 5",
    size=22, color=ICE_BLUE, font=F_BODY,
)

# Pie autor
add_text_box(
    s, Inches(0.8), Inches(6.4), Inches(10), Inches(0.4),
    "YOLOv8 · NCNN · Edge Computing",
    size=12, color=ICE_BLUE,
)
add_text_box(
    s, Inches(0.8), Inches(6.8), Inches(10), Inches(0.4),
    "Suiza · 2026",
    size=12, bold=True, color=WHITE,
)


# ───────── SLIDE 2: Agenda ─────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_page_title(s, "Agenda", "Recorrido por la propuesta")

# Dos columnas de bloques numerados
agenda_items = [
    ("01", "El problema",           "Por qué detectar armas automáticamente"),
    ("02", "La solución",            "Arquitectura y enfoque general"),
    ("03", "Tecnologías clave",      "Python, YOLOv8, Pi5, NCNN"),
    ("04", "Cómo funciona",          "Pipeline paso a paso"),
    ("05", "Entrenamiento del modelo","Dataset y proceso"),
    ("06", "Anti-alucinaciones",     "Reducir falsos positivos"),
    ("07", "Métricas y resultados",  "Lo que conseguimos"),
    ("08", "Despliegue y futuro",    "Pi5 + roadmap"),
]

col_w = Inches(5.8)
row_h = Inches(0.7)
start_top = Inches(1.7)
gap = Inches(0.15)

for i, (num, title, desc) in enumerate(agenda_items):
    col = i % 2
    row = i // 2
    left = Inches(0.65 + col * 6.2)
    top  = start_top + (row_h + gap) * row

    # Número grande
    add_text_box(s, left, top, Inches(1.0), row_h,
                 num, size=32, bold=True, color=ACCENT, font=F_TITLE,
                 anchor=MSO_ANCHOR.MIDDLE)
    # Texto
    add_text_box(s, left + Inches(1.0), top + Inches(0.05), col_w - Inches(1.0), Inches(0.35),
                 title, size=16, bold=True, color=NAVY)
    add_text_box(s, left + Inches(1.0), top + Inches(0.38), col_w - Inches(1.0), Inches(0.3),
                 desc, size=11, color=MUTED)

add_footer(s, 2, TOTAL_SLIDES)


# ───────── SLIDE 3: El problema ─────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_page_title(s, "El problema", "Limitaciones de la vigilancia humana")

# Tres columnas con icono + título + descripción
problems = [
    ("👁", "Cobertura limitada",
     "Un operador no puede mirar 20 pantallas a la vez con la misma atención sostenida."),
    ("⏱", "Fatiga visual",
     "Tras una hora de turno, la capacidad de detectar anomalías cae drásticamente."),
    ("⚡", "Tiempo de reacción",
     "Entre detectar y avisar pueden pasar segundos críticos en una emergencia."),
]
col_w = Inches(3.9)
gap = Inches(0.25)
total_w = col_w * 3 + gap * 2
start_left = (prs.slide_width - total_w) // 2
top = Inches(2.0)

for i, (icon, title, desc) in enumerate(problems):
    left = start_left + (col_w + gap) * i
    # Tarjeta
    card = add_filled_rect(s, left, top, col_w, Inches(3.6), LIGHT_BG)
    # Banda superior accent
    add_filled_rect(s, left, top, col_w, Inches(0.18), ACCENT)
    # Icono grande
    add_text_box(s, left, top + Inches(0.5), col_w, Inches(1.0),
                 icon, size=60, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Título
    add_text_box(s, left + Inches(0.3), top + Inches(1.7), col_w - Inches(0.6), Inches(0.5),
                 title, size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    # Descripción
    add_text_box(s, left + Inches(0.3), top + Inches(2.2), col_w - Inches(0.6), Inches(1.2),
                 desc, size=12, color=DARK_TXT, align=PP_ALIGN.CENTER, line_spacing=1.25)

# Mensaje cierre
add_text_box(
    s, Inches(0.65), Inches(6.0), Inches(12), Inches(0.5),
    "La IA no reemplaza al humano — le da una ventaja: nunca se distrae y reacciona en milisegundos.",
    size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
)

add_footer(s, 3, TOTAL_SLIDES)


# ───────── SLIDE 4: Objetivos ─────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_page_title(s, "Objetivos del proyecto", "General y específicos")

# Objetivo general (caja grande izquierda)
add_filled_rect(s, Inches(0.65), Inches(1.7), Inches(5.5), Inches(4.5), NAVY)
add_text_box(s, Inches(0.95), Inches(1.95), Inches(4.9), Inches(0.5),
             "OBJETIVO GENERAL", size=12, bold=True, color=ICE_BLUE)
add_text_box(s, Inches(0.95), Inches(2.5), Inches(4.9), Inches(3.5),
             "Construir un sistema de visión por computadora que detecte armas en tiempo real, "
             "con alta fiabilidad, ejecutándose íntegramente en hardware de bajo coste "
             "(Raspberry Pi 5).",
             size=16, color=WHITE, line_spacing=1.35)

# Objetivos específicos (columna derecha con bullets)
add_text_box(s, Inches(6.6), Inches(1.7), Inches(6.5), Inches(0.4),
             "OBJETIVOS ESPECÍFICOS", size=12, bold=True, color=ACCENT)

specifics = [
    "Detectar 3 categorías: pistola, arma larga y arma blanca.",
    "Reconocer armas en distintos ángulos (lado, frontal, cenital).",
    "Minimizar falsos positivos sobre objetos confundibles.",
    "Funcionar a 5–10 FPS en Raspberry Pi 5.",
    "Procesamiento 100% local — sin internet, sin nube.",
]
top = Inches(2.2)
for i, item in enumerate(specifics):
    # punto coloreado
    add_filled_rect(s, Inches(6.6), top + Inches(0.15), Inches(0.12), Inches(0.12), ACCENT)
    add_text_box(s, Inches(6.85), top, Inches(6.3), Inches(0.7),
                 item, size=14, color=DARK_TXT, line_spacing=1.25)
    top += Inches(0.75)

add_footer(s, 4, TOTAL_SLIDES)


# ───────── SLIDE 5: La solución ─────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_page_title(s, "La solución", "Tres bloques que trabajan en conjunto")

blocks = [
    ("A", "Captura",
     "Cámara USB o módulo CSI conectada a la Pi5. Captura video continuo."),
    ("B", "Inferencia",
     "Modelo YOLOv8 nano analiza cada frame en ~100 ms y devuelve detecciones."),
    ("C", "Respuesta",
     "Si la confianza supera el umbral, alerta inmediata: guarda evidencia y notifica."),
]
col_w = Inches(3.9)
gap = Inches(0.25)
total_w = col_w * 3 + gap * 2
start_left = (prs.slide_width - total_w) // 2
top = Inches(1.9)

for i, (letter, title, desc) in enumerate(blocks):
    left = start_left + (col_w + gap) * i
    # Letra grande
    add_text_box(s, left, top, col_w, Inches(2.0),
                 letter, size=120, bold=True, color=ICE_BLUE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=F_TITLE)
    # Título
    add_text_box(s, left, top + Inches(1.85), col_w, Inches(0.5),
                 title, size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    # Descripción
    add_text_box(s, left + Inches(0.3), top + Inches(2.5), col_w - Inches(0.6), Inches(1.5),
                 desc, size=13, color=DARK_TXT, align=PP_ALIGN.CENTER, line_spacing=1.3)
    # Flecha entre bloques
    if i < 2:
        arrow_left = left + col_w + Inches(0.02)
        add_text_box(s, arrow_left, top + Inches(0.7), gap, Inches(0.6),
                     "→", size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Cita destacada
add_filled_rect(s, Inches(0.65), Inches(6.0), Inches(12), Inches(0.85), LIGHT_BG)
add_text_box(s, Inches(0.95), Inches(6.15), Inches(11.4), Inches(0.6),
             "Todo el procesamiento es LOCAL en la Pi5. Sin nube, sin internet, sin riesgo de privacidad.",
             size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 5, TOTAL_SLIDES)


# ───────── SLIDE 6: Tecnologías ─────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_page_title(s, "Tecnologías utilizadas", "Stack open source de la industria")

techs = [
    ("Python",          "Lenguaje principal — el estándar en IA."),
    ("YOLOv8",          "Modelo de detección en tiempo real, una sola pasada."),
    ("Ultralytics",     "Librería para entrenar/inferir/exportar YOLOv8."),
    ("PyTorch",         "Motor matemático bajo YOLO — operaciones de la red neuronal."),
    ("OpenCV",          "Captura de cámara, dibujo de cajas, guardado de frames."),
    ("ONNX + NCNN",     "Formatos optimizados. NCNN ejecuta en Pi5 2-3× más rápido."),
]
cols = 3
rows = 2
col_w = Inches(4.0)
row_h = Inches(2.2)
gap_x = Inches(0.18)
gap_y = Inches(0.22)
total_w = col_w * cols + gap_x * (cols - 1)
start_left = (prs.slide_width - total_w) // 2
top = Inches(1.8)

for i, (name, desc) in enumerate(techs):
    col = i % cols
    row = i // cols
    left = start_left + (col_w + gap_x) * col
    cell_top = top + (row_h + gap_y) * row
    # Borde lateral accent
    add_filled_rect(s, left, cell_top, Inches(0.12), row_h, ACCENT)
    # Fondo light
    add_filled_rect(s, left + Inches(0.12), cell_top, col_w - Inches(0.12), row_h, LIGHT_BG)
    # Texto
    add_text_box(s, left + Inches(0.35), cell_top + Inches(0.25), col_w - Inches(0.55), Inches(0.55),
                 name, size=20, bold=True, color=NAVY)
    add_text_box(s, left + Inches(0.35), cell_top + Inches(0.95), col_w - Inches(0.55), Inches(1.1),
                 desc, size=12, color=DARK_TXT, line_spacing=1.3)

add_footer(s, 6, TOTAL_SLIDES)


# ───────── SLIDE 7: Hardware Pi5 ─────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_page_title(s, "¿Por qué Raspberry Pi 5?", "Hardware modesto, resultados profesionales")

# Tabla comparativa estilo cards
specs = [
    ("80 €",        "Coste",            "vs ~1.000 € PC con GPU"),
    ("~10 W",       "Consumo",          "vs ~80 W PC dedicado"),
    ("5–7 FPS",     "Velocidad real",   "en inferencia con YOLOv8n"),
    ("Local",       "Privacidad",       "todo el procesamiento on-device"),
]
col_w = Inches(2.95)
gap = Inches(0.18)
total_w = col_w * 4 + gap * 3
start_left = (prs.slide_width - total_w) // 2
top = Inches(2.1)

for i, (big, label, sub) in enumerate(specs):
    left = start_left + (col_w + gap) * i
    add_filled_rect(s, left, top, col_w, Inches(3.6), NAVY)
    # Gran número
    add_text_box(s, left, top + Inches(0.5), col_w, Inches(1.4),
                 big, size=48, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=F_TITLE)
    # Label
    add_text_box(s, left, top + Inches(2.0), col_w, Inches(0.5),
                 label.upper(), size=14, bold=True, color=ICE_BLUE,
                 align=PP_ALIGN.CENTER)
    # Sub
    add_text_box(s, left + Inches(0.25), top + Inches(2.6), col_w - Inches(0.5), Inches(0.8),
                 sub, size=11, color=ICE_BLUE, align=PP_ALIGN.CENTER, line_spacing=1.3)

# Nota inferior
add_text_box(
    s, Inches(0.65), Inches(6.1), Inches(12), Inches(0.5),
    "Procesador ARM Cortex-A76 de 4 núcleos a 2.4 GHz · 4 u 8 GB de RAM",
    size=13, color=MUTED, align=PP_ALIGN.CENTER,
)

add_footer(s, 7, TOTAL_SLIDES)


# ───────── SLIDE 8: Arquitectura / Pipeline ─────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_page_title(s, "Arquitectura del sistema", "El flujo de datos, de la cámara a la alerta")

# Pipeline horizontal con 5 cajas
steps = [
    ("📷", "Cámara",       "Video 30 FPS"),
    ("🔧", "Preprocesa",    "Redim. 416×416"),
    ("🧠", "YOLOv8n",       "Inferencia ~100ms"),
    ("⚙",  "Filtro NMS",   "conf ≥ 0.5"),
    ("🚨", "Alerta",        "Caja + log"),
]
n = len(steps)
box_w = Inches(2.05)
arrow_w = Inches(0.45)
total_w = box_w * n + arrow_w * (n - 1)
start_left = (prs.slide_width - total_w) // 2
top = Inches(2.2)

for i, (icon, name, sub) in enumerate(steps):
    left = start_left + (box_w + arrow_w) * i
    # Caja
    add_filled_rect(s, left, top, box_w, Inches(2.6), LIGHT_BG)
    add_filled_rect(s, left, top, box_w, Inches(0.5), NAVY)
    # número arriba
    add_text_box(s, left, top + Inches(0.05), box_w, Inches(0.4),
                 f"PASO {i+1}", size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # icono
    add_text_box(s, left, top + Inches(0.65), box_w, Inches(0.9),
                 icon, size=46, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # nombre
    add_text_box(s, left, top + Inches(1.6), box_w, Inches(0.4),
                 name, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    # sub
    add_text_box(s, left, top + Inches(2.0), box_w, Inches(0.5),
                 sub, size=11, color=MUTED, align=PP_ALIGN.CENTER, line_spacing=1.2)
    # flecha
    if i < n - 1:
        arrow_left = left + box_w + Inches(0.02)
        add_text_box(s, arrow_left, top + Inches(0.9), arrow_w, Inches(0.6),
                     "▶", size=24, bold=True, color=ACCENT,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Pie clave
add_text_box(s, Inches(0.65), Inches(5.5), Inches(12), Inches(0.5),
             "Cada frame se procesa de forma independiente — el sistema no necesita memoria del pasado.",
             size=13, color=DARK_TXT, align=PP_ALIGN.CENTER)
add_text_box(s, Inches(0.65), Inches(6.0), Inches(12), Inches(0.5),
             "Total: ~150-200 ms por frame  →  5-7 FPS reales en Pi5",
             size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_footer(s, 8, TOTAL_SLIDES)


# ───────── SLIDE 9: Cómo funciona paso a paso ─────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_page_title(s, "Cómo funciona, paso a paso", "Qué hace el sistema cuando aparece un arma")

steps5 = [
    ("1", "Captura de frame",
     "La cámara entrega una imagen a la Pi5 a 30 FPS. OpenCV la lee como matriz de píxeles."),
    ("2", "Preprocesamiento",
     "Se redimensiona a 416×416 y se normalizan los valores de píxel (0–1)."),
    ("3", "Inferencia",
     "YOLOv8n analiza la imagen y devuelve cajas, clases y confianzas en ~100 ms."),
    ("4", "Filtros",
     "Se descartan detecciones débiles (conf < 0.5) y cajas duplicadas (NMS)."),
    ("5", "Respuesta",
     "Se dibuja la caja, se guarda el frame de evidencia y se registra la alerta."),
]

col_w = Inches(11.8)
top = Inches(1.7)
row_h = Inches(0.95)
gap = Inches(0.08)

for i, (num, title, desc) in enumerate(steps5):
    y = top + (row_h + gap) * i
    # Número grande en círculo
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.65), y + Inches(0.05),
                                  Inches(0.8), Inches(0.8))
    circle.fill.solid(); circle.fill.fore_color.rgb = NAVY
    circle.line.fill.background()
    circle.shadow.inherit = False
    tf = circle.text_frame
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE
    r.font.name = F_TITLE

    # Texto
    add_text_box(s, Inches(1.65), y, col_w - Inches(1.0), Inches(0.4),
                 title, size=16, bold=True, color=NAVY)
    add_text_box(s, Inches(1.65), y + Inches(0.42), col_w - Inches(1.0), Inches(0.5),
                 desc, size=12, color=DARK_TXT, line_spacing=1.3)

add_footer(s, 9, TOTAL_SLIDES)


# ───────── SLIDE 10: Entrenamiento del modelo ─────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_page_title(s, "El modelo de IA: cómo se entrena", "Aprendizaje supervisado con miles de ejemplos")

# Izquierda: explicación
add_text_box(s, Inches(0.65), Inches(1.85), Inches(6.5), Inches(0.5),
             "Aprende por ejemplos", size=18, bold=True, color=NAVY)
add_text_box(s, Inches(0.65), Inches(2.4), Inches(6.5), Inches(3.5),
             "Al modelo se le muestran miles de imágenes etiquetadas con "
             "cajas indicando dónde está cada arma.\n\n"
             "Durante 30 épocas, ajusta millones de pesos internos para "
             "minimizar el error entre lo que predice y la etiqueta real.\n\n"
             "Tras el entrenamiento, generaliza: detecta armas en imágenes "
             "que nunca vio.",
             size=13, color=DARK_TXT, line_spacing=1.4)

# Derecha: stats key
right_left = Inches(7.6)
right_top = Inches(1.85)
add_filled_rect(s, right_left, right_top, Inches(5.1), Inches(4.5), NAVY)

stats = [
    ("35.842",    "imágenes en el dataset"),
    ("3",         "categorías (knife / handgun / long_gun)"),
    ("30",        "épocas de entrenamiento"),
    ("~1.5 h",    "entrenando el modelo nano (RTX 4060)"),
    ("3 M / 11 M","parámetros nano / small"),
]
y = right_top + Inches(0.35)
for big, label in stats:
    add_text_box(s, right_left + Inches(0.3), y, Inches(2.2), Inches(0.55),
                 big, size=24, bold=True, color=WHITE, font=F_TITLE)
    add_text_box(s, right_left + Inches(2.6), y + Inches(0.1), Inches(2.4), Inches(0.45),
                 label, size=11, color=ICE_BLUE, line_spacing=1.2)
    y += Inches(0.83)

add_footer(s, 10, TOTAL_SLIDES)


# ───────── SLIDE 11: Anti-alucinaciones ─────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_page_title(s, "Estrategia anti-alucinaciones", "Evitar que el modelo confunda objetos cotidianos con armas")

# Texto introductorio
add_filled_rect(s, Inches(0.65), Inches(1.75), Inches(12), Inches(0.8), LIGHT_BG)
add_text_box(
    s, Inches(0.95), Inches(1.85), Inches(11.4), Inches(0.6),
    "Falso positivo = el modelo cree ver un arma donde hay un móvil, mando, paraguas o bolsa. En un stand, eso destruiría la fiabilidad.",
    size=12, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
)

# Tres tácticas
tactics = [
    ("Hard negatives",
     "2.500 imágenes de objetos parecidos PERO sin armas se añaden al entrenamiento."),
    ("Penalización extra",
     "El componente cls de la loss se sube de 0.5 a 1.0 — castiga el doble los falsos positivos."),
    ("Umbral alto",
     "En despliegue se exige conf ≥ 0.5 (en lugar del 0.25 por defecto)."),
]
col_w = Inches(4.0)
gap = Inches(0.15)
total_w = col_w * 3 + gap * 2
start_left = (prs.slide_width - total_w) // 2
top = Inches(2.95)

for i, (name, desc) in enumerate(tactics):
    left = start_left + (col_w + gap) * i
    # Card
    add_filled_rect(s, left, top, col_w, Inches(2.8), LIGHT_BG)
    add_filled_rect(s, left, top, col_w, Inches(0.5), ACCENT)
    add_text_box(s, left, top + Inches(0.05), col_w, Inches(0.4),
                 f"TÁCTICA {i+1}", size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(s, left + Inches(0.25), top + Inches(0.75), col_w - Inches(0.5), Inches(0.6),
                 name, size=17, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text_box(s, left + Inches(0.3), top + Inches(1.4), col_w - Inches(0.6), Inches(1.3),
                 desc, size=12, color=DARK_TXT, align=PP_ALIGN.CENTER, line_spacing=1.3)

# Resultado
add_text_box(s, Inches(0.65), Inches(6.05), Inches(12), Inches(0.55),
             "Resultado: precisión del nano +2.4 puntos porcentuales — muchas menos falsas alarmas.",
             size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_footer(s, 11, TOTAL_SLIDES)


# ───────── SLIDE 12: Métricas ─────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_page_title(s, "Métricas de rendimiento", "Lo que el modelo realmente consigue")

# Tabla de resultados
headers = ["Modelo", "Parámetros", "mAP@50", "Precisión", "Recall", "Plataforma"]
rows = [
    ["Nano v4",  "3 M",  "0.852", "0.875", "0.758", "Raspberry Pi 5"],
    ["Small v4", "11 M", "0.879", "0.879", "0.792", "PC con GPU"],
]

cell_w = [Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5), Inches(2.5)]
table_left = Inches(1.4)
top = Inches(1.9)
row_h = Inches(0.7)

# Header
left = table_left
for j, h in enumerate(headers):
    add_filled_rect(s, left, top, cell_w[j], row_h, NAVY)
    add_text_box(s, left, top, cell_w[j], row_h,
                 h, size=13, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    left += cell_w[j]

# Rows
for i, r in enumerate(rows):
    y = top + row_h * (i + 1)
    left = table_left
    bg = LIGHT_BG if i % 2 == 0 else WHITE
    for j, c in enumerate(r):
        add_filled_rect(s, left, y, cell_w[j], row_h, bg)
        is_metric = j in (2, 3, 4)
        add_text_box(s, left, y, cell_w[j], row_h,
                     c, size=14, bold=(j == 0) or is_metric,
                     color=NAVY if (j == 0) else (ACCENT if is_metric else DARK_TXT),
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        left += cell_w[j]

# Interpretación grande
add_text_box(s, Inches(0.65), Inches(4.2), Inches(12), Inches(0.5),
             "Lectura clave del nano (Pi5):", size=16, bold=True, color=NAVY)

reads = [
    ("87.5 %", "de las veces que dice \"arma\" acierta"),
    ("75.8 %", "de las armas reales son detectadas"),
    ("5-7 FPS","velocidad en Raspberry Pi 5"),
]
col_w = Inches(4.0)
gap = Inches(0.15)
total_w = col_w * 3 + gap * 2
start_left = (prs.slide_width - total_w) // 2
read_top = Inches(4.85)

for i, (big, sub) in enumerate(reads):
    left = start_left + (col_w + gap) * i
    add_filled_rect(s, left, read_top, col_w, Inches(1.7), LIGHT_BG)
    add_text_box(s, left, read_top + Inches(0.15), col_w, Inches(0.9),
                 big, size=38, bold=True, color=ACCENT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=F_TITLE)
    add_text_box(s, left + Inches(0.3), read_top + Inches(1.1), col_w - Inches(0.6), Inches(0.5),
                 sub, size=12, color=DARK_TXT, align=PP_ALIGN.CENTER, line_spacing=1.2)

add_footer(s, 12, TOTAL_SLIDES)


# ───────── SLIDE 13: Despliegue ─────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_page_title(s, "Despliegue en Raspberry Pi 5", "Del modelo entrenado a la cámara en producción")

# Izquierda: pasos
add_text_box(s, Inches(0.65), Inches(1.85), Inches(6.5), Inches(0.5),
             "Pasos de despliegue", size=18, bold=True, color=NAVY)

deploy_steps = [
    "Exportar el modelo nano a NCNN (formato optimizado ARM)",
    "Copiar la carpeta (~12 MB) a la Pi5 vía SCP",
    "Instalar dependencias en la Pi5 (ultralytics + ncnn)",
    "Conectar cámara USB / módulo CSI",
    "Lanzar el script de inferencia con conf=0.5",
]
y = Inches(2.5)
for i, st in enumerate(deploy_steps):
    # Número en círculo accent
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.65), y + Inches(0.05),
                                Inches(0.5), Inches(0.5))
    circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    circle.shadow.inherit = False
    tf = circle.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(i + 1)
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE
    r.font.name = F_TITLE
    # Texto
    add_text_box(s, Inches(1.35), y + Inches(0.05), Inches(5.8), Inches(0.5),
                 st, size=12, color=DARK_TXT, line_spacing=1.3)
    y += Inches(0.65)

# Derecha: bloque comando
add_filled_rect(s, Inches(7.6), Inches(1.85), Inches(5.1), Inches(4.5), NAVY)
add_text_box(s, Inches(7.9), Inches(2.05), Inches(4.6), Inches(0.5),
             "COMANDO EN LA Pi5", size=11, bold=True, color=ICE_BLUE)

add_text_box(s, Inches(7.9), Inches(2.6), Inches(4.6), Inches(2.5),
             "yolo predict \\\n  model=best_ncnn_model \\\n  imgsz=416 \\\n  source=0 \\\n  conf=0.5",
             size=14, color=WHITE, font="Consolas", line_spacing=1.35)

add_text_box(s, Inches(7.9), Inches(5.3), Inches(4.6), Inches(1),
             "source=0 → cámara conectada\nconf=0.5 → umbral anti-alucinaciones",
             size=11, color=ICE_BLUE, line_spacing=1.4)

add_footer(s, 13, TOTAL_SLIDES)


# ───────── SLIDE 14: Limitaciones ─────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_page_title(s, "Limitaciones honestas", "Lo que el sistema NO puede hacer")

limits = [
    ("Armas a gran distancia",
     "Si la pistola ocupa menos de 20×20 píxeles, la detección no es fiable."),
    ("Oclusiones parciales",
     "Un arma oculta bajo chaqueta o dentro del bolsillo no es visible."),
    ("Iluminación extrema",
     "Escenas muy oscuras o con contraluces fuertes degradan el rendimiento."),
    ("Latencia",
     "Si el arma aparece y desaparece en <100 ms, puede no detectarse."),
    ("Armas no convencionales",
     "El modelo no reconoce armas caseras, replicas raras o disfraces."),
    ("No sustituye al humano",
     "Es una capa de apoyo al operador, no un reemplazo."),
]
cols = 2
col_w = Inches(6.0)
gap_x = Inches(0.2)
row_h = Inches(1.4)
gap_y = Inches(0.15)
total_w = col_w * cols + gap_x
start_left = (prs.slide_width - total_w) // 2
top = Inches(1.85)

for i, (title, desc) in enumerate(limits):
    col = i % cols
    row = i // cols
    left = start_left + (col_w + gap_x) * col
    cell_top = top + (row_h + gap_y) * row
    add_filled_rect(s, left, cell_top, col_w, row_h, LIGHT_BG)
    add_filled_rect(s, left, cell_top, Inches(0.1), row_h, RGBColor(0xD8, 0x4A, 0x4A))
    add_text_box(s, left + Inches(0.3), cell_top + Inches(0.2), col_w - Inches(0.5), Inches(0.4),
                 title, size=15, bold=True, color=NAVY)
    add_text_box(s, left + Inches(0.3), cell_top + Inches(0.65), col_w - Inches(0.5), Inches(0.7),
                 desc, size=11, color=DARK_TXT, line_spacing=1.3)

add_footer(s, 14, TOTAL_SLIDES)


# ───────── SLIDE 15: Próximos pasos ─────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, WHITE)
add_page_title(s, "Próximos pasos", "Hoja de ruta tras esta propuesta")

nexts = [
    ("Adaptación al escenario real",
     "Capturar 50-100 frames del stand y reentrenar para ajustar al contexto específico."),
    ("Tracking entre frames",
     "Confirmar detecciones con persistencia temporal — reduce alertas momentáneas."),
    ("Integración con notificaciones",
     "Conectar con Telegram, email o panel web para alertas remotas."),
    ("Cuantización INT8",
     "Acelerar aún más la inferencia en Pi5 con pérdida mínima de precisión."),
    ("Reconocimiento de comportamiento",
     "No solo el objeto: apuntar, blandir, ocultar — análisis de pose y acción."),
]
top = Inches(1.85)
row_h = Inches(0.9)
gap = Inches(0.12)

for i, (title, desc) in enumerate(nexts):
    y = top + (row_h + gap) * i
    # Caja accent left
    add_filled_rect(s, Inches(0.65), y, Inches(0.1), row_h, ACCENT)
    add_text_box(s, Inches(0.95), y + Inches(0.05), Inches(11), Inches(0.4),
                 title, size=15, bold=True, color=NAVY)
    add_text_box(s, Inches(0.95), y + Inches(0.45), Inches(11), Inches(0.45),
                 desc, size=12, color=DARK_TXT, line_spacing=1.3)

add_footer(s, 15, TOTAL_SLIDES)


# ───────── SLIDE 16: Conclusión / Cierre ─────────
s = prs.slides.add_slide(BLANK)
set_slide_bg(s, NAVY)
add_filled_rect(s, Inches(0), Inches(0), Inches(0.3), Inches(7.5), ACCENT)

add_text_box(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "EN RESUMEN", size=14, bold=True, color=ICE_BLUE)

add_text_box(s, Inches(0.8), Inches(1.75), Inches(11.5), Inches(2.0),
             "Detección de armas con IA,\nfiable y al alcance.",
             size=44, bold=True, color=WHITE, line_spacing=1.1, font=F_TITLE)

# Tres pilares
pillars = [
    ("Fiabilidad",     "87.5% precisión, anti-alucinaciones activas"),
    ("Accesibilidad",  "Corre en hardware de 80 €"),
    ("Privacidad",     "100% local, sin nube"),
]
col_w = Inches(4.0)
gap = Inches(0.15)
total_w = col_w * 3 + gap * 2
start_left = (prs.slide_width - total_w) // 2
top = Inches(4.5)

for i, (h, sub) in enumerate(pillars):
    left = start_left + (col_w + gap) * i
    add_text_box(s, left, top, col_w, Inches(0.6),
                 h, size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(s, left + Inches(0.2), top + Inches(0.7), col_w - Inches(0.4), Inches(1.0),
                 sub, size=13, color=ICE_BLUE, align=PP_ALIGN.CENTER, line_spacing=1.3)

# Cierre
add_text_box(s, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
             "Gracias.", size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.3),
             "¿Preguntas?",
             size=14, color=ICE_BLUE, align=PP_ALIGN.CENTER)


# ───────── Guardar ─────────
prs.save(str(OUT))
print(f"PPTX generado: {OUT}")
print(f"Slides: {len(prs.slides)}")
